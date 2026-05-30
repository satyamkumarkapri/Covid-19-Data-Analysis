from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import datetime

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "COVID-19 Data Analysis Dashboard"
    subtitle.text = f"Data Science & Analytics Project\n{datetime.datetime.now().strftime('%Y-%m-%d')}"

    # Slide 2: Project Overview
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Project Overview"
    tf = body_shape.text_frame
    tf.text = "A full-stack Web Application & Data Science Dashboard."
    p = tf.add_paragraph()
    p.text = "Visualizes and analyzes COVID-19 pandemic data."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Uses Glassmorphism UI, Flask backend, and Machine Learning."
    p.level = 1

    # Slide 3: Workflow & Tech Stack
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Workflow & Tech Stack"
    tf = body_shape.text_frame
    tf.text = "Frontend: HTML5, CSS3, JavaScript, Bootstrap 5"
    p = tf.add_paragraph()
    p.text = "Backend: Python, Flask"
    p = tf.add_paragraph()
    p.text = "Data Science: Pandas, NumPy, Scikit-Learn, SciPy"
    p = tf.add_paragraph()
    p.text = "Methodology: M1 to M5 (Cleaning, EDA, Prob, Stats, Regression)"

    # Slide 4: Dataset Explanation
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Dataset Explanation"
    tf = body_shape.text_frame
    tf.text = "Source: Johns Hopkins University CSSE"
    p = tf.add_paragraph()
    p.text = "time_series_covid19_confirmed_global.csv"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "time_series_covid19_deaths_global.csv"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Data covers hundreds of countries dynamically over years."

    # Slide 5: Dashboard & Charts
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Dashboard Features & Charts"
    tf = body_shape.text_frame
    tf.text = "Interactive Chart.js and Plotly graphs:"
    p = tf.add_paragraph()
    p.text = "Global Trend Line & Histogram"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Donut/Pie charts for status distribution"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Interactive Country Data Table with Sorting & Searching"
    p.level = 1

    # Slide 6: ML Prediction
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Machine Learning Prediction"
    tf = body_shape.text_frame
    tf.text = "Algorithm: Linear Regression"
    p = tf.add_paragraph()
    p.text = "Feature: Time elapsed (Days since outbreak)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Target: Total Confirmed Cases"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Output: 30-day continuous future projection."

    # Slide 7: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Conclusion & Future Scope"
    tf = body_shape.text_frame
    tf.text = "Conclusion: Successfully combined web-dev and data science."
    p = tf.add_paragraph()
    p.text = "Future Scope:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Integration of Live API Feeds"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Deep Learning (LSTM) for non-linear predictions"
    p.level = 1

    prs.save('presentation.pptx')
    print("presentation.pptx successfully generated.")

if __name__ == '__main__':
    create_presentation()
